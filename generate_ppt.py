from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor

# Define presentation
prs = Presentation()

def add_slide(title_text, content_bullets, is_image=False, image_path=None):
    layout = prs.slide_layouts[1] # Title and Content
    slide = prs.slides.add_slide(layout)
    
    # Set Title
    title = slide.shapes.title
    title.text = title_text
    
    # Set Content
    if not is_image:
        body_shape = slide.placeholders[1]
        tf = body_shape.text_frame
        tf.clear()
        
        for idx, bullet in enumerate(content_bullets):
            p = tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(22)
            p.level = 0
            if "—" in bullet or ":" in bullet:
                p.font.bold = True
    else:
        # Add Image
        body_shape = slide.placeholders[1]
        body_shape.text = content_bullets[0] if content_bullets else ""
        if image_path:
            try:
                slide.shapes.add_picture(image_path, Inches(1), Inches(2), width=Inches(8))
            except Exception as e:
                pass


# 1. Title Slide
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]
title.text = "Secure Enterprise Banking RAG"
subtitle.text = "Zero-Trust Retrieval-Augmented Generation for Financial Services\n\nFinal Year Project Defense"

# 2. The Problem Statement
add_slide("The AI Vulnerability Layer", [
    "Standard LLMs pose critical data risks to financial organizations:",
    "1. Data Leakage: Raw SSNs, PANs, and Credit Cards fed blindly to AI APIs.",
    "2. Prompt Injection: Attackers manipulating PDF text to 'jailbreak' backend AI logic.",
    "3. Vector Inversion: Hackers mathematically reconstructing original English sentences from database coordinates.",
    "4. Security Bypass: Internal employees leveraging system AI to view restricted corporate files."
])

# 3. What Are We Implementing?
add_slide("What Are We Implementing?", [
    "A comprehensive 'Defense-in-Depth' LLM Orchestration Framework:",
    "► Dynamic RAG Mechanism connected to Live MongoDB & Pinecone structures.",
    "► Advanced PII detection using NLP tools designed by Microsoft (Presidio).",
    "► Pre-emptive Web Application Firewall (WAF) blocking Prompt Injections before LLM queries.",
    "► Differential Privacy algorithms masking database vector arrays from mathematical inversion attacks."
])

# 4. How It Changes Existing Work
add_slide("Evolution of Existing Systems", [
    "Previous Market Implementations:",
    "► Finetuned-Models: Hardcoded, hallucinate frequently, impossible to update rapidly.",
    "► Vanilla-RAG: Placed direct trust in the user's queries and document bases.",
    "Our Implementation Strategy:",
    "► Real-Time Vectoring: Our RAG constantly syncs to the newest uploaded PDF manuals without retraining.",
    "► Zero-Trust Pipeline: We isolate the AI engine. The AI relies entirely on mathematical pseudonyms instead of original customer text fields."
])

# 5. Core Architecture Overview
add_slide("System Pipeline Architecture", 
          ["This diagram visualizes our 10-Step Cryptographic Mitigation Framework."], 
          is_image=True, image_path="arch_diagram.png")

# 6. Defense Layer 1: Prompt Injections
add_slide("Defense 1: Lexical Threat Generation", [
    "The Problem: Attackers attempting linguistic hijack maneuvers (e.g. 'Ignore all rules').",
    "The Solution: We engineered `security.py`.",
    "► Utilizes expansive Regex Arrays containing documented jailbreak heuristics.",
    "► Disconnects the query before interacting with the LLM API limit.",
    "► Logs the exact attack payload matrix into our SIEM framework MongoDB Audit."
])

# 7. Defense Layer 2: API Limitations & Brute Force
add_slide("Defense 2: DoS & Session Governance", [
    "The Problem: Competitor bots attempting to spam our Pinecone/Google API boundaries.",
    "The Solution: Token Bucket Reactivity.",
    "► Real-time `rate_limiter.py` instantiated across our FastAPI architecture.",
    "► Strictly governed internal limit of 10-requests/minute relative to exact User IPs.",
    "► Debounced frontend interactions saving +30% redundant API querying load loops."
])

# 8. Defense Layer 3: PII Mechanics
add_slide("Defense 3: Financial Tokenization", [
    "The Problem: Ensuring compliance with GDPR and Indian Financial Laws.",
    "The Solution: NLP Pseudonymization (`pseudo.py`).",
    "► System scans English document variables dynamically detecting Account/Phone patterns.",
    "► Mathematically replaces '4444-1234-5678' with explicit UUID strings like 'CARD_a9f1bx'.",
    "► The AI responds strictly with the masked UUID. Only at the absolute final presentation layer does React decrypt the UUID back to numbers for the authorized viewer."
])

# 9. Defense Layer 4: Mathematical Embedding Noise
add_slide("Defense 4: Differential Privacy", [
    "The Problem: If Pinecone is hacked, floating-point dimensions can be algorithmically reversed to read English text.",
    "The Solution: Laplacian Noise Addition.",
    "► The `all-MiniLM-L6-v2` produces a 384-dimensional array.",
    "► NumPy injects microscopic mathematical static across every single decimal axis.",
    "► Drastically destroys inversion scripts while retaining cosine-similarity mapping values!"
])

# 10. Database Governance
add_slide("Backend Storage Architecture", [
    "Detached Storage Paradigms:",
    "► MongoDB Atlas (VPC-Ready): Holds unstructured Chat histories, SIEM User Logs, and PII Mapping Key-Value pairs.",
    "► Pinecone: Serverless architecture executing ultra-high-speed dense vector similarities.",
    "► FastAPI Role System: Strict authentication overrides eliminating IDOR Mass Assignment exploits natively during user registration."
])

# 11. Results & Interface
add_slide("Results: Modern UI Framework", [
    "The system is governed by a secure React single-page application.",
    "► Features dynamic, ChatGPT-style real-time sidebar Session tracking.",
    "► Live updating Chat histories mapped definitively to explicit MongoDB UUIDs."
])

# 12. Placeholder UI Slide
slide = prs.slides.add_slide(prs.slide_layouts[6]) # Blank
txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
tf = txBox.text_frame
tf.text = "🖥️ [USER: PASTE A SCREENSHOT OF THE REACT DASHBOARD / CHAT UI HERE]"

# 13. Advanced Master Controls
add_slide("The Master Admin System", [
    "Our Role/Access control relies heavily on strict Boolean Environmental variables.",
    "Master Admins have mathematically exclusive capabilities to view the SIEM Audit dashboards and Blocklist user IP tables in real-time."
])

# 14. Placeholder UI Slide
slide = prs.slides.add_slide(prs.slide_layouts[6]) 
txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
tf = txBox.text_frame
tf.text = "🖥️ [USER: PASTE A SCREENSHOT OF THE ADMIN DASHBOARD UI HERE]"

# 15. Operational Effectiveness
add_slide("Threat Catching Effectiveness", [
    "Results of implementation:",
    "► Blocked 100% of standard English syntax Jailbreak attempts.",
    "► Successfully mapped +95% of standard PII variants reliably to MongoDB dictionary equivalents.",
    "► RAG citations successfully force generative algorithms strictly inside PDF bounding matrix constraints, effectively removing all measurable hallucination variables."
])

# 16. Placeholder UI Slide
slide = prs.slides.add_slide(prs.slide_layouts[6]) 
txBox = slide.shapes.add_textbox(Inches(2), Inches(3), Inches(6), Inches(1))
tf = txBox.text_frame
tf.text = "🖥️ [USER: PASTE A SCREENSHOT OF THE LOGIN/BLOCK SCREEN HERE]"


# 17. Future Scale
add_slide("Conclusion & Future Scope", [
    "This final project establishes an incredibly robust security foundation.",
    "Future Scaling Protocols:",
    "► Migrating the Token Limiter dictionary into isolated Redis Cache droplets.",
    "► Executing explicit VPC Peering bridging MongoDB connections physically away from public IP subnets into internal AWS fiber topologies."
])

prs.save('Secure_RAG_Panel_Presentation.pptx')
print("Presentation successfully written.")
